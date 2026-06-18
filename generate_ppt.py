import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    c_dark_blue  = RGBColor(15, 23, 42)     # #0F172A
    c_muted_gray = RGBColor(100, 116, 139)  # #64748B
    c_light_bg   = RGBColor(248, 250, 252)  # #F8FAFC
    c_border_gray= RGBColor(226, 232, 240)  # #E2E8F0
    
    # Section Colors
    # 1. Slate/Grey
    h_grey   = RGBColor(51, 65, 85)       # #334155
    bg_grey  = RGBColor(241, 245, 249)    # #F1F5F9
    b_grey   = RGBColor(203, 213, 225)    # #CBD5E1
    
    # 2. Crimson Red
    h_red    = RGBColor(153, 27, 27)      # #991B1B
    bg_red   = RGBColor(254, 242, 242)     # #FEF2F2
    b_red    = RGBColor(252, 165, 165)    # #FCA5A5
    
    # 3. Violet/Purple
    h_purple = RGBColor(109, 40, 217)     # #6D28D9
    bg_purple= RGBColor(245, 243, 255)    # #F5F3FF
    b_purple = RGBColor(221, 214, 254)    # #DDD6FE

    # -------------------------------------------------------------
    # Template Helper: Applies slide background and standard headers/footers
    # -------------------------------------------------------------
    def apply_slide_template(slide, title_text, subtitle_text):
        # Subtle light gray background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = c_light_bg
        bg_shape.line.fill.background()
        
        # Vertical sidebar decorations
        sidebar_w = Inches(0.25)
        left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sidebar_w, prs.slide_height)
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = c_dark_blue
        left_bar.line.fill.background()
        
        right_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, prs.slide_width - sidebar_w, 0, sidebar_w, prs.slide_height)
        right_bar.fill.solid()
        right_bar.fill.fore_color.rgb = c_dark_blue
        right_bar.line.fill.background()
        
        # Bottom divider line
        bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.15), prs.slide_width - Inches(1.0), Inches(0.02))
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = c_border_gray
        bottom_line.line.fill.background()

        # Bottom footer text
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.20), Inches(6.0), Inches(0.25))
        p_f = footer_box.text_frame.paragraphs[0]
        p_f.text = "RiskIntel: AI-Powered Bankruptcy Intelligence Portal  |  Client Demo Presentation"
        p_f.font.name = "Calibri"
        p_f.font.size = Pt(7.5)
        p_f.font.color.rgb = c_muted_gray

        copyright_box = slide.shapes.add_textbox(prs.slide_width - Inches(3.0), Inches(7.20), Inches(2.5), Inches(0.25))
        p_c = copyright_box.text_frame.paragraphs[0]
        p_c.text = "© 2026 RiskIntel. Confidential."
        p_c.font.name = "Calibri"
        p_c.font.size = Pt(7.5)
        p_c.font.color.rgb = c_muted_gray

        # Main Title and Subtitle
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11.0), Inches(0.7))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = "Calibri"
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = c_dark_blue
        
        p_sub = tf_t.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(10.0)
        p_sub.font.italic = True
        p_sub.font.color.rgb = c_muted_gray

    # -------------------------------------------------------------
    # Content Card Helper (for bullet lists)
    # -------------------------------------------------------------
    def add_card(slide, title, bullets, left, top, width, height, header_color, body_bg_color, border_color):
        # Draw body rectangle
        body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        body.fill.solid()
        body.fill.fore_color.rgb = body_bg_color
        body.line.color.rgb = border_color
        body.line.width = Pt(1.2)
        
        # Draw header rectangle
        h_height = Inches(0.38)
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h_height)
        header.fill.solid()
        header.fill.fore_color.rgb = header_color
        header.line.fill.background()
        
        # Header text
        tf_h = header.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = Inches(0.12)
        tf_h.margin_right = Inches(0.12)
        tf_h.margin_top = Inches(0.06)
        tf_h.margin_bottom = Inches(0.06)
        
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.name = "Calibri"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = RGBColor(255, 255, 255)
        
        # Body text box
        tb = slide.shapes.add_textbox(left, top + h_height, width, height - h_height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        
        for idx, item in enumerate(bullets):
            text_str, level = item if isinstance(item, tuple) else (item, 0)
            
            p = tf.add_paragraph() if idx > 0 or len(tf.paragraphs[0].text) > 0 else tf.paragraphs[0]
            p.level = level
            p.space_after = Pt(3.0)
            
            # Format bold prefix if colon exists
            if ":" in text_str and level == 0:
                parts = text_str.split(":", 1)
                prefix = parts[0] + ":"
                suffix = parts[1]
                
                run1 = p.add_run()
                run1.text = prefix
                run1.font.bold = True
                run1.font.size = Pt(9.0)
                run1.font.name = "Calibri"
                run1.font.color.rgb = RGBColor(15, 23, 42)
                
                run2 = p.add_run()
                run2.text = suffix
                run2.font.bold = False
                run2.font.size = Pt(9.0)
                run2.font.name = "Calibri"
                run2.font.color.rgb = RGBColor(51, 65, 85)
            else:
                run = p.add_run()
                run.text = text_str
                run.font.bold = False
                run.font.size = Pt(9.0 if level == 0 else 8.5)
                run.font.name = "Calibri"
                run.font.color.rgb = RGBColor(51, 65, 85)

    # Use blank layout for all slides
    blank_layout = prs.slide_layouts[6]

    # =============================================================================
    # SLIDE 1: Business Objective & Problem Statement
    # =============================================================================
    slide_1 = prs.slides.add_slide(blank_layout)
    apply_slide_template(
        slide_1, 
        "Business Objective & Problem Statement", 
        "Establishing the Operational Framework and Business Rationale for Automated Bankruptcy Operations"
    )
    
    bullets_problem = [
        "Inefficient Manual Discovery: Monitoring bankruptcy filings across 94 US court districts relies on labor-intensive daily tracking, introducing high latency and critical operational gaps.",
        "Recovery Value Leakage: Delays in recognizing chapter conversions (e.g. Ch 13 payment plan converting to Ch 7 liquidation) result in missed Proof of Claim (POC) filing deadlines and outright credit write-offs.",
        "Data Silos & Query Backlogs: Business risk managers cannot directly interrogate legal portfolios, generating long queues for database administrators and data analysts.",
        "Inconsistent Docket Standards: Court feeds contain varied formats, abbreviations, and missing fields (such as Pro Se indicators or trustee assignments) that complicate automated workflows."
    ]
    add_card(
        slide_1, "The Problem Statement", bullets_problem,
        Inches(0.5), Inches(1.25), Inches(5.9), Inches(5.6),
        h_red, bg_red, b_red
    )

    bullets_objective = [
        "Real-Time Exposure Visibility: Automate the processing of bankruptcy notices to identify risk accounts in near real-time, reducing latency from days to under an hour.",
        "Asset Preservation & Claim Capture: Standardize the detection of chapter transitions, ensuring automated triggers alert risk teams to immediately file Proofs of Claim (POC) and freeze credit limits.",
        "Self-Service Natural Language Analytics: Empower credit officers and legal managers with an AI-driven interface to run ad-hoc compliance audits without writing complex SQL.",
        "Cohesive Legal Workbench: Consolidate filings, chapter conversions, pro se indicators, defense counsels, and court districts into a single client-facing workspace."
    ]
    add_card(
        slide_1, "The Business Objective", bullets_objective,
        Inches(6.8), Inches(1.25), Inches(5.9), Inches(5.6),
        h_grey, bg_grey, b_grey
    )


    # =============================================================================
    # SLIDE 2: Potential Business Challenges to Solve
    # =============================================================================
    slide_2 = prs.slides.add_slide(blank_layout)
    apply_slide_template(
        slide_2,
        "Potential Business Challenges to Solve",
        "Addressing Key Operational Risk, Compliance Regulations, and Technical Bottlenecks"
    )

    bullets_op_challenges = [
        "Docket Explosion & Volume: Handling hundreds of court filings daily without scalable text processing systems causes review backlogs.",
        "Trustee & Attorney Attribution: Inconsistent naming conventions make it difficult to group cases by defense firms or trustees.",
        "Resource Dependencies: Risk managers are bottlenecked by specialized SQL developers whenever custom dashboard metrics are needed."
    ]
    add_card(
        slide_2, "Operational Challenges", bullets_op_challenges,
        Inches(0.5), Inches(1.25), Inches(3.9), Inches(5.6),
        h_purple, bg_purple, b_purple
    )

    bullets_rc_challenges = [
        "Strict Court Deadlines: Proof of Claim (POC) filing deadlines are non-negotiable; missing them permanently forfeits recovery assets.",
        "Self-Represented (Pro Se) Debtors: Debtors without legal counsel exhibit unpredictable filing behavior and high compliance complexity.",
        "Line-of-Credit Exposure: Delays in freezing credit cards or active loans after a filing allows debtors to draw down remaining funds."
    ]
    add_card(
        slide_2, "Risk & Compliance Challenges", bullets_rc_challenges,
        Inches(4.7), Inches(1.25), Inches(3.9), Inches(5.6),
        h_red, bg_red, b_red
    )

    bullets_tech_challenges = [
        "Legacy Registry Feeds: Integrating fragmented, non-standardized court database schema mappings.",
        "LLM Query Accuracy: Preventing LLM 'hallucinations' in SQL code generation to ensure queries execute safely on transaction servers.",
        "API Scalability & Cost: Guarding against compounding LLM invocation costs under repeated operational inquiries."
    ]
    add_card(
        slide_2, "Technical & Data Challenges", bullets_tech_challenges,
        Inches(8.9), Inches(1.25), Inches(3.9), Inches(5.6),
        h_grey, bg_grey, b_grey
    )


    # =============================================================================
    # SLIDE 3: Solution Overview (RiskIntel Overview Dashboard)
    # =============================================================================
    slide_3 = prs.slides.add_slide(blank_layout)
    apply_slide_template(
        slide_3,
        "RiskIntel Overview Dashboard: Systems Design",
        "AI-Driven Bankruptcy Monitoring, Dynamic Query INTERFACE, and Portfolio Intelligence Workspace"
    )

    # Left Column
    col1_left  = Inches(0.5)
    col1_width = Inches(3.9)
    
    bullets_challenges_s3 = [
        "Data Complexity: High-volume bankruptcy filings across multiple court jurisdictions with heterogeneous schemas.",
        "Manual Audit Overhead: Heavy manual effort required to monitor court dockets, status dates, and creditor meetings.",
        "Invisible Trailing Risks: Undetected chapter conversions and silent pro se filings leading to operational losses.",
        "Integration Fractures: Disconnect between loan account registers, credit databases, and live court data streams."
    ]
    add_card(
        slide_3, "Business Challenges", bullets_challenges_s3,
        col1_left, Inches(1.25), col1_width, Inches(3.1),
        h_grey, bg_grey, b_grey
    )
    
    bullets_coverage_s3 = [
        "Automated Multi-State Coverage: Californian, Texan, Floridian, and national court registers scanned daily.",
        "Core Chapters Monitored: Continuous status checking on Chapter 7, 11, and 13 cases.",
        "Conversion Pipelines Tracked: Automated notifications on original chapter updates (e.g. Ch 13 -> 7).",
        "Key Actor Demographics: Categorized directories of Pro Se parties, defense firms, and presiding trustees."
    ]
    add_card(
        slide_3, "Portfolio & Case Coverage", bullets_coverage_s3,
        col1_left, Inches(4.55), col1_width, Inches(2.3),
        h_red, bg_red, b_red
    )
    
    # Middle Column
    col2_left  = Inches(4.6)
    col2_width = Inches(3.9)
    
    bullets_features_s3 = [
        ("AI-Powered SQL Generator:", 0),
        ("  • Translates business user questions to executable SQL queries.", 0),
        ("  • Performs context validation to prevent database syntax errors.", 0),
        ("Autonomous Trend Detection:", 0),
        ("  • Automatically surfaces high-volume peaks and geo-concentrations.", 0),
        ("  • Identifies significant changes in filing distribution patterns.", 0),
        ("Predictive Risk Modeling:", 0),
        ("  • Generates customized Match Scores for active case risk profiles.", 0),
        ("  • Flags high-risk debtor accounts before standard reporting lag.", 0),
        ("Adaptive Data Visualization:", 0),
        ("  • Auto-scales chart layout and dimensions based on query scale.", 0),
        ("  • Generates count-enriched pie wedges and clean bar charts.", 0),
        ("Interactive Analytics Tabs:", 0),
        ("  • Seamless drill-downs for attorney activity, asset types, and timelines.", 0)
    ]
    add_card(
        slide_3, "Salient Features", bullets_features_s3,
        col2_left, Inches(1.25), col2_width, Inches(5.6),
        h_purple, bg_purple, b_purple
    )
    
    # Right Column
    col3_left  = Inches(8.7)
    col3_width = Inches(4.1)
    
    bullets_solution_s3 = [
        "Central Risk Workbench: A Streamlit interface providing unified database visibility for operational managers.",
        "Dual-LLM Engine: Orchestrates AWS Bedrock (Claude/Llama) and OpenAI endpoints to support generation and validation.",
        "Smart Cache Layer: Employs LangChain and SQLite caching to deliver sub-second response times and zero token costs.",
        "Decoupled Architecture: Separates core ingestion, querying, forecasting, and visualization modules for resilience."
    ]
    add_card(
        slide_3, "Solution and Approach", bullets_solution_s3,
        col3_left, Inches(1.25), col3_width, Inches(2.4),
        h_grey, bg_grey, b_grey
    )
    
    bullets_benefits_s3 = [
        "Portfolio Visibility: Over 80% reduction in query cycle time using natural language instead of manual reports.",
        "Proactive Mitigation: Automated conversion alerts reduce losses by flagging active asset liquidations early.",
        "Operational Productivity: Minimizes dependency on specialized SQL developers to write custom filters.",
        "Decision Accuracy: Multi-LLM validation check ensures generated analytics are contextually precise."
    ]
    add_card(
        slide_3, "Benefits", bullets_benefits_s3,
        col3_left, Inches(3.85), col3_width, Inches(2.0),
        h_red, bg_red, b_red
    )
    
    bullets_tech_s3 = [
        "Python | Streamlit | AWS Bedrock | OpenAI | SQLite | Pandas | NumPy | Matplotlib | Seaborn | Plotly | LangChain"
    ]
    add_card(
        slide_3, "Technologies", bullets_tech_s3,
        col3_left, Inches(6.05), col3_width, Inches(0.8),
        h_purple, bg_purple, b_purple
    )


    # =============================================================================
    # SLIDE 4: Cloud Platform Architecture Solutions (AWS Swim-Lane Diagram)
    # =============================================================================
    slide_4 = prs.slides.add_slide(blank_layout)
    apply_slide_template(
        slide_4,
        "Cloud Platform Architecture Solution",
        "Secure, Enterprise-Grade Ingestion Pipeline and LLM Inference Flow on AWS Cloud"
    )

    # ── Color palette for this slide ───────────────────────────────────────────
    aws_orange   = RGBColor(255, 153,   0)   # AWS brand orange
    aws_blue     = RGBColor( 35,  47,  62)   # AWS dark navy
    zone_ingst   = RGBColor(229, 243, 255)   # Light blue  – Ingestion zone
    zone_store   = RGBColor(232, 246, 235)   # Light green – Storage zone
    zone_ai      = RGBColor(243, 232, 255)   # Light purple – AI/ML zone
    zone_portal  = RGBColor(255, 243, 229)   # Light amber – Portal zone
    arrow_color  = RGBColor(100, 116, 139)   # Slate arrow colour

    # ── Icon path helper ───────────────────────────────────────────────────────
    import os
    base_dir = os.path.dirname(os.path.abspath(__file__))
    def icon(name):
        p = os.path.join(base_dir, name)
        return p if os.path.exists(p) else None

    # ── Layout constants ───────────────────────────────────────────────────────
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    SWIM_TOP    = Inches(1.30)   # top of swim-lane rows
    SWIM_H      = Inches(2.35)   # height of each swim-lane row
    SWIM_GAP    = Inches(0.12)   # gap between rows
    NODE_W      = Inches(1.55)   # node box width
    NODE_H      = Inches(1.30)   # node box height
    ICON_SZ     = Inches(0.38)   # icon square size
    HDR_H       = Inches(0.28)   # header bar height

    # ── 5 horizontal zones (columns) ──────────────────────────────────────────
    #  Zone 0: Data Sources  (outside AWS)
    #  Zone 1: Ingestion     (AWS Lambda / API GW)
    #  Zone 2: Storage       (S3 / Cache)
    #  Zone 3: AI / ML       (Bedrock + SQL Judge)
    #  Zone 4: Delivery      (ECS / Portal / Users)

    ZONE_LABELS = ["Data Sources", "Ingestion Layer", "Storage Layer", "AI / ML Layer", "Delivery Layer"]
    ZONE_COLORS = [RGBColor(241,245,249), zone_ingst, zone_store, zone_ai, zone_portal]
    ZONE_BORDER = [RGBColor(203,213,225), RGBColor(147,197,253), RGBColor(134,239,172),
                   RGBColor(196,181,253), RGBColor(253,186,116)]

    # Compute zone lefts/widths across the full slide (minus 0.5" padding each side)
    CONTENT_LEFT  = Inches(0.38)
    CONTENT_RIGHT = slide_w - Inches(0.38)
    CONTENT_W     = CONTENT_RIGHT - CONTENT_LEFT
    zone_w = CONTENT_W / 5

    # Total height: two swim lanes + gap
    TOTAL_SWIM_H = SWIM_H * 2 + SWIM_GAP
    SWIM_BOT = SWIM_TOP + TOTAL_SWIM_H  # bottom edge

    # ── Draw zone background panels ───────────────────────────────────────────
    for zi, (zlabel, zcol, zborder) in enumerate(zip(ZONE_LABELS, ZONE_COLORS, ZONE_BORDER)):
        zl = CONTENT_LEFT + zone_w * zi
        zbox = slide_4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, zl, SWIM_TOP, zone_w - Inches(0.06), TOTAL_SWIM_H)
        zbox.fill.solid()
        zbox.fill.fore_color.rgb = zcol
        zbox.line.color.rgb = zborder
        zbox.line.width = Pt(1.0)

        # Zone label (top of zone)
        ztb = slide_4.shapes.add_textbox(zl + Inches(0.05), SWIM_TOP + Inches(0.05), zone_w - Inches(0.15), Inches(0.28))
        p_z = ztb.text_frame.paragraphs[0]
        p_z.text = zlabel
        p_z.font.name = "Calibri"
        p_z.font.size = Pt(8.5)
        p_z.font.bold = True
        p_z.font.color.rgb = aws_blue if zi > 0 else RGBColor(71, 85, 105)
        p_z.alignment = PP_ALIGN.CENTER

    # ── Node helper ───────────────────────────────────────────────────────────
    def draw_node(slide, label, sublabel, left, top, w, h, hdr_color, body_color, border_color, icon_path=None):
        # Body
        body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        body.fill.solid()
        body.fill.fore_color.rgb = body_color
        body.line.color.rgb = border_color
        body.line.width = Pt(1.2)

        # Header bar
        hb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, HDR_H)
        hb.fill.solid()
        hb.fill.fore_color.rgb = hdr_color
        hb.line.fill.background()

        # Icon
        icon_offset = Inches(0.0)
        if icon_path:
            try:
                slide.shapes.add_picture(icon_path, left + Inches(0.05), top + HDR_H + Inches(0.06),
                                          ICON_SZ, ICON_SZ)
                icon_offset = ICON_SZ + Inches(0.05)
            except Exception as e:
                print(f"  [warn] icon {icon_path}: {e}")

        # Header text
        ht = hb.text_frame
        ht.word_wrap = True
        ht.margin_left = Inches(0.07)
        ht.margin_top = Inches(0.02)
        p_ht = ht.paragraphs[0]
        p_ht.text = label
        p_ht.font.name = "Calibri"
        p_ht.font.size = Pt(8.0)
        p_ht.font.bold = True
        p_ht.font.color.rgb = RGBColor(255, 255, 255)
        p_ht.alignment = PP_ALIGN.CENTER

        # Sub-label text (body)
        btb = slide.shapes.add_textbox(left + icon_offset, top + HDR_H + Inches(0.04),
                                        w - icon_offset - Inches(0.04), h - HDR_H - Inches(0.08))
        btf = btb.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.04)
        btf.margin_right = Inches(0.04)
        btf.margin_top = Inches(0.04)
        p_bt = btf.paragraphs[0]
        p_bt.text = sublabel
        p_bt.font.name = "Calibri"
        p_bt.font.size = Pt(7.0)
        p_bt.font.color.rgb = RGBColor(51, 65, 85)
        p_bt.alignment = PP_ALIGN.CENTER

    # ── Arrow helpers ─────────────────────────────────────────────────────────
    def arrow_r(slide, left, top, width, label=""):
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.18))
        arr.fill.solid(); arr.fill.fore_color.rgb = arrow_color; arr.line.fill.background()
        if label:
            atb = slide.shapes.add_textbox(left, top - Inches(0.18), width, Inches(0.18))
            p = atb.text_frame.paragraphs[0]
            p.text = label; p.font.name = "Calibri"; p.font.size = Pt(6.5)
            p.font.color.rgb = c_muted_gray; p.alignment = PP_ALIGN.CENTER

    def arrow_d(slide, left, top, height, label=""):
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.18), height)
        arr.fill.solid(); arr.fill.fore_color.rgb = arrow_color; arr.line.fill.background()
        if label:
            atb = slide.shapes.add_textbox(left + Inches(0.22), top + height / 2 - Inches(0.1), Inches(0.9), Inches(0.2))
            p = atb.text_frame.paragraphs[0]
            p.text = label; p.font.name = "Calibri"; p.font.size = Pt(6.5)
            p.font.color.rgb = c_muted_gray

    def arrow_u(slide, left, top, height, label=""):
        arr = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, left, top, Inches(0.18), height)
        arr.fill.solid(); arr.fill.fore_color.rgb = arrow_color; arr.line.fill.background()

    def arrow_l(slide, left, top, width, label=""):
        arr = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, left, top, width, Inches(0.18))
        arr.fill.solid(); arr.fill.fore_color.rgb = arrow_color; arr.line.fill.background()

    # ── Compute node positions ─────────────────────────────────────────────────
    # Nodes centre-aligned inside each zone
    # Row A (top swim lane): Court API, Lambda, S3, SQL Judge, ECS / Portal
    # Row B (bottom swim lane): Account DB, API GW, Cache DB, Bedrock, Risk Users

    def zone_cx(zi):   # centre-x of zone i
        return CONTENT_LEFT + zone_w * zi + zone_w / 2

    def node_left(zi):
        return zone_cx(zi) - NODE_W / 2

    ROW_A_TOP = SWIM_TOP + Inches(0.42)   # top of row A nodes
    ROW_B_TOP = SWIM_TOP + SWIM_H + SWIM_GAP + Inches(0.42)  # top of row B nodes

    # ── Row A nodes ───────────────────────────────────────────────────────────
    # Zone 0 – External Court APIs
    draw_node(slide_4, "Court Data Feeds", "Live PACER dockets\nJSON / XML court notices",
              node_left(0), ROW_A_TOP, NODE_W, NODE_H,
              h_red, RGBColor(255,255,255), b_red)

    # Zone 1 – AWS Lambda
    draw_node(slide_4, "AWS Lambda", "Event-driven ingestion\n& schema validation",
              node_left(1), ROW_A_TOP, NODE_W, NODE_H,
              h_grey, RGBColor(255,255,255), b_grey,
              icon("aws_lambda_icon.png"))

    # Zone 2 – Amazon S3
    draw_node(slide_4, "Amazon S3", "KMS-encrypted raw\n& processed dockets",
              node_left(2), ROW_A_TOP, NODE_W, NODE_H,
              h_purple, RGBColor(255,255,255), b_purple,
              icon("aws_s3_icon.png"))

    # Zone 3 – SQL Judge
    draw_node(slide_4, "SQL Safety Judge", "Validates & sanitizes\nLLM-generated queries",
              node_left(3), ROW_A_TOP, NODE_W, NODE_H,
              h_red, RGBColor(255,255,255), b_red)

    # Zone 4 – ECS Fargate / Analytics
    draw_node(slide_4, "ECS Fargate Engine", "Query execution &\nvisualization engine",
              node_left(4), ROW_A_TOP, NODE_W, NODE_H,
              h_grey, RGBColor(255,255,255), b_grey,
              icon("aws_ecs_icon.png"))

    # ── Row B nodes ───────────────────────────────────────────────────────────
    # Zone 0 – Account Registers
    draw_node(slide_4, "Portfolio Registers", "Internal credit &\naccount master data",
              node_left(0), ROW_B_TOP, NODE_W, NODE_H,
              h_red, RGBColor(255,255,255), b_red)

    # Zone 1 – API Gateway
    draw_node(slide_4, "AWS API Gateway", "Authenticated REST\nendpoints & rate limiting",
              node_left(1), ROW_B_TOP, NODE_W, NODE_H,
              h_grey, RGBColor(255,255,255), b_grey)

    # Zone 2 – Cache DB
    draw_node(slide_4, "SQLite Cache DB", "LangChain prompt cache\n& schema store",
              node_left(2), ROW_B_TOP, NODE_W, NODE_H,
              h_purple, RGBColor(255,255,255), b_purple)

    # Zone 3 – Bedrock LLM
    draw_node(slide_4, "AWS Bedrock LLM", "Claude-3.5 / Llama\nmanaged inference",
              node_left(3), ROW_B_TOP, NODE_W, NODE_H,
              h_purple, RGBColor(255,255,255), b_purple,
              icon("aws_bedrock_icon.png"))

    # Zone 4 – RiskIntel Portal + Users
    draw_node(slide_4, "RiskIntel Portal", "Streamlit UI\nNL query & dashboards",
              node_left(4), ROW_B_TOP, NODE_W, NODE_H,
              h_grey, RGBColor(255,255,255), b_grey)

    # ── Flow arrows ───────────────────────────────────────────────────────────
    # Arrow mid-Y helpers
    def mid_a(row_top): return row_top + NODE_H / 2 - Inches(0.09)   # midpoint of node row

    # Gap between zones (horizontal space between node right edge and next zone node left edge)
    gap_h = zone_cx(1) - NODE_W / 2 - (zone_cx(0) + NODE_W / 2)

    def zone_right(zi, row_top):
        return node_left(zi) + NODE_W

    def arrow_between_zones(zi_from, zi_to, row_top, label=""):
        x_start = node_left(zi_from) + NODE_W + Inches(0.03)
        x_end   = node_left(zi_to)   - Inches(0.03)
        arrow_r(slide_4, x_start, mid_a(row_top), x_end - x_start, label)

    # Row A horizontal flow: Court -> Lambda -> S3 -> SQL Judge -> ECS
    arrow_between_zones(0, 1, ROW_A_TOP, "feeds")
    arrow_between_zones(1, 2, ROW_A_TOP, "raw events")
    arrow_between_zones(2, 3, ROW_A_TOP, "validated SQL")
    arrow_between_zones(3, 4, ROW_A_TOP, "safe query")

    # Row B horizontal flow: Portfolio -> API GW -> Cache -> Bedrock -> Portal
    arrow_between_zones(0, 1, ROW_B_TOP, "account data")
    arrow_between_zones(1, 2, ROW_B_TOP, "schema lookup")
    arrow_between_zones(2, 3, ROW_B_TOP, "NL query")
    arrow_between_zones(3, 4, ROW_B_TOP, "results")

    # Vertical arrows between rows A and B (within same zone)
    def v_arrow_top_to_bot(zi):
        cx = zone_cx(zi) - Inches(0.09)
        y_start = ROW_A_TOP + NODE_H + Inches(0.04)
        y_end   = ROW_B_TOP - Inches(0.04)
        arrow_d(slide_4, cx, y_start, y_end - y_start)

    def v_arrow_bot_to_top(zi):
        cx = zone_cx(zi) + Inches(0.04)
        y_start = ROW_A_TOP + NODE_H + Inches(0.04)
        y_end   = ROW_B_TOP - Inches(0.04)
        arrow_u(slide_4, cx, y_start, y_end - y_start)

    # S3 feeds down into Cache (zone 2)
    v_arrow_top_to_bot(2)
    # SQL Judge fed from Bedrock (zone 3) — up from Bedrock to Judge
    v_arrow_bot_to_top(3)
    # ECS result -> Portal (zone 4)
    v_arrow_top_to_bot(4)

    # ── Legend box ────────────────────────────────────────────────────────────
    legend_top = SWIM_BOT + Inches(0.12)
    legend_items = [
        (h_red,    "External / Security"),
        (h_grey,   "Compute / Runtime"),
        (h_purple, "Storage / AI-ML"),
    ]
    lx = CONTENT_LEFT + Inches(0.2)
    for lcolor, llabel in legend_items:
        dot = slide_4.shapes.add_shape(MSO_SHAPE.OVAL, lx, legend_top + Inches(0.04), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = lcolor; dot.line.fill.background()
        ltb = slide_4.shapes.add_textbox(lx + Inches(0.18), legend_top, Inches(1.5), Inches(0.22))
        p_l = ltb.text_frame.paragraphs[0]
        p_l.text = llabel; p_l.font.name = "Calibri"; p_l.font.size = Pt(7.5)
        p_l.font.color.rgb = c_muted_gray
        lx += Inches(1.75)

    # Save the presentation
    output_filename = "RiskIntel_Overview_V5.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully to: {output_filename}")

if __name__ == "__main__":
    create_presentation()
